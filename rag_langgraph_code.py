import os
from typing import List, Dict, Any
from dataclasses import dataclass, field
from langgraph.graph import StateGraph, END
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.vectorstores import FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.document_loaders import PyPDFLoader, TextLoader
from langchain.docstore.document import Document
from tqdm import tqdm
from transformers import AutoModelForQuestionAnswering, AutoTokenizer, pipeline
from docx import Document as DocxDocument
from pptx import Presentation

@dataclass
class GraphState:
    query: str = ""
    documents: List[Document] = field(default_factory=list)
    answer: str = ""

def load_docx(file_path: str) -> List[Document]:
    docx = DocxDocument(file_path)
    text = "\n".join([paragraph.text for paragraph in docx.paragraphs if paragraph.text.strip()])
    return [Document(page_content=text)]

def load_ppt(file_path: str) -> List[Document]:
    presentation = Presentation(file_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text += paragraph.text + "\n"
    return [Document(page_content=text)]

def load_documents(file_path: str) -> List[Document]:
    print("Loading document...")
    
    file_extension = file_path.lower().split('.')[-1]
    
    if file_extension == 'pdf':
        loader = PyPDFLoader(file_path)
        documents = loader.load()
    elif file_extension == 'txt':
        loader = TextLoader(file_path, encoding="utf-8")
        documents = loader.load()
    elif file_extension == 'docx':
        documents = load_docx(file_path)
    elif file_extension == 'pptx':
        documents = load_ppt(file_path)
    else:
        raise ValueError(f"Unsupported file format: {file_extension}")

    print("Splitting documents...")
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=100
    )
    split_docs = list(tqdm(text_splitter.split_documents(documents), desc="Splitting"))
    return split_docs

def setup_vector_store(documents: List[Document]):
    embedding_model = HuggingFaceEmbeddings(
        model_name="sentence-transformers/all-mpnet-base-v2"
    )
    print("Creating vector store...")
    docs_with_progress = list(tqdm(documents, desc="Embedding"))
    vector_store = FAISS.from_documents(
        documents=docs_with_progress,
        embedding=embedding_model
    )
    return vector_store, embedding_model

def create_qa_pipeline():
    model_name = "deepset/roberta-base-squad2"  
    model = AutoModelForQuestionAnswering.from_pretrained(model_name)
    tokenizer = AutoTokenizer.from_pretrained(model_name)
    qa_pipeline = pipeline("question-answering", model=model, tokenizer=tokenizer)
    return qa_pipeline

def create_rag_workflow(vector_store):
    reader = create_qa_pipeline()

    def retrieve(state: GraphState) -> Dict[str, Any]:
        documents = vector_store.similarity_search(state.query, k=3)
        return {"documents": documents}

    def generate_answer(state: GraphState) -> Dict[str, Any]:
        documents = state.documents
        query = state.query

        context = "\n".join([doc.page_content for doc in documents])
        qa_input = {"question": query, "context": context}

        try:
            answer = reader(qa_input)["answer"]
        except Exception as e:
            answer = f"Error generating answer: {str(e)}"

        return {"answer": answer}

    workflow = StateGraph(GraphState)
    workflow.add_node("retrieve", retrieve)
    workflow.add_node("generate", generate_answer)
    workflow.set_entry_point("retrieve")
    workflow.add_edge("retrieve", "generate")
    workflow.add_edge("generate", END)

    return workflow.compile()

def main(file_path: str):
    documents = load_documents(file_path)
    vector_store, _ = setup_vector_store(documents)
    app = create_rag_workflow(vector_store)
    while True:
        query = input("\nEnter your query (or type 'exit' to quit): ").strip()
        
        if query.lower() == 'exit':
            break
        initial_state = GraphState(query=query)
        result = app.invoke(initial_state)

        print("\nQuery Result:")
        print(f"Answer: {result['answer']}")

if __name__ == "__main__":
    file_path = input("Enter the path to your document (PDF/TXT/DOCX/PPTX): ").strip()
    main(file_path)